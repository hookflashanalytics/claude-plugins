"""
Populate the Hookflash 5-slide Results Analysis template.

Requires python-pptx (`pip install python-pptx`) and the 5-slide template
(downloaded via Tether `get_slide_template(template="results_analysis")`, or
bundled). The template is FIXED at 5 slides; this fills them, nothing more:

  Slide 1 — fill the text placeholders (experiment name; hypothesis IF/THEN/
            BECAUSE; test type; pages; audience; primary + secondary metrics) and
            add a small overall conversion-rate chart on the right.
  Slide 2 — fill the table with OVERALL scores (one row per variation).
  Slide 3 — big overall conversion-rate bar chart.
  Slide 4 — fill the table with DEVICE-SPLIT scores (variation × device).
  Slide 5 — big device-split conversion-rate bar chart.

Usage:  python build_deck.py <template.pptx> <results.json> <out.pptx>
  results.json = {"summary": {experiment, test_type, pages, audience,
                    primary_metric, secondary_metrics, hypothesis_if,
                    hypothesis_then, hypothesis_because},
                  "results": { <the tapa_ra results object> }}
"""
from __future__ import annotations
import json, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Site accent palette (cyan / magenta / green). CONFIRM the exact brand hexes.
ACCENTS = ["22D3EE", "D96AF6", "22E06B"]
DEVICE_ORDER = ["overall", "mobile", "desktop", "tablet"]


def _pct(x):
    return "—" if x is None else f"{x * 100:.2f}%"


def _signed(x):
    return "—" if x is None else f"{(x * 100):+.2f}%"


def _conf(x):
    return "—" if x is None else f"{x * 100:.1f}%"


# ---------- slide 1: text ----------
def _iter_runs(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                yield run


def replace_everywhere(prs, old, new):
    for slide in prs.slides:
        for run in _iter_runs(slide):
            if old in run.text:
                run.text = run.text.replace(old, new)


def fill_slide1_placeholders(slide, summary):
    """Replace the ordered 'Insert Text' placeholders on slide 1. Order matches
    the template: IF, THEN, BECAUSE, test type, pages, audience, primary metric,
    secondary metric(s)."""
    values = [
        summary.get("hypothesis_if", ""),
        summary.get("hypothesis_then", ""),
        summary.get("hypothesis_because", ""),
        summary.get("test_type", ""),
        summary.get("pages", ""),
        summary.get("audience", ""),
        summary.get("primary_metric", ""),
        summary.get("secondary_metrics", ""),
    ]
    i = 0
    for run in _iter_runs(slide):
        if "Insert Text" in run.text and i < len(values):
            run.text = run.text.replace("Insert Text", values[i] or "—", 1)
            i += 1


# ---------- tables ----------
def _first_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def _set_row(table, r, cells):
    for c, val in enumerate(cells):
        if c < len(table.columns) and r < len(table.rows):
            table.cell(r, c).text = str(val)


def fill_overall_table(slide, variations):
    """Overall slide table has NO Device column:
    Variation | Users | Conversions | Conv. rate | Uplift | Confidence.
    Row 0 is the header; one data row per variation. The control (first
    variation) has blank uplift/confidence."""
    table = _first_table(slide)
    if table is None:
        return
    for idx, v in enumerate(variations):
        _set_row(table, idx + 1, [
            v["name"], f"{v['users']:,}", f"{v['conversions']:,}",
            _pct(v["rate"]), _signed(v.get("uplift")), _conf(v.get("confidence")),
        ])


def fill_device_table(slide, variations):
    """One row per (variation × device) in DEVICE_ORDER, skipping devices absent
    from the data. Row 0 is the header."""
    table = _first_table(slide)
    if table is None:
        return
    r = 1
    for v in variations:
        by_dev = v.get("by_device", {})
        for dev in DEVICE_ORDER:
            cell = {"users": v["users"], "conversions": v["conversions"], "rate": v["rate"]} \
                if dev == "overall" else by_dev.get(dev)
            if not cell:
                continue
            _set_row(table, r, [
                v["name"], dev.capitalize(), f"{cell['users']:,}", f"{cell['conversions']:,}",
                _pct(cell["rate"]),
                _signed(v.get("uplift")) if dev == "overall" else "",
                _conf(v.get("confidence")) if dev == "overall" else "",
            ])
            r += 1


# ---------- charts ----------
def add_bar_chart(slide, categories, values, title, x, y, cx, cy):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Conversion rate %", values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(x), Inches(y), Inches(cx), Inches(cy), data)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    plot = chart.plots[0]
    plot.vary_by_categories = True
    # brand-colour each bar from the accent palette
    points = plot.series[0].points
    for i, pt in enumerate(points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = RGBColor.from_string(ACCENTS[i % len(ACCENTS)])
    return chart


def overall_series(variations):
    return [v["name"] for v in variations], [round(v["rate"] * 100, 2) for v in variations]


def device_categories(variations):
    """Grouped-by-device categories: for a single-series chart we flatten to
    'Variation — Device' with its conversion rate."""
    cats, vals = [], []
    for v in variations:
        for dev in DEVICE_ORDER[1:]:  # skip Overall for the device chart
            cell = v.get("by_device", {}).get(dev)
            if cell:
                cats.append(f"{v['name']} — {dev.capitalize()}")
                vals.append(round(cell["rate"] * 100, 2))
    return cats, vals


def build(template_path, payload, out_path):
    prs = Presentation(template_path)
    summary = payload.get("summary", {})
    results = payload["results"]
    variations = results["kpis"][0]["variations"]  # primary KPI drives the deck
    slides = list(prs.slides)

    # Slide 1
    if summary.get("experiment"):
        replace_everywhere(prs, "Insert Experiment Name", summary["experiment"])
    fill_slide1_placeholders(slides[0], summary)
    cats, vals = overall_series(variations)
    add_bar_chart(slides[0], cats, vals, "Overall conversion rate", 7.6, 1.9, 5.2, 4.2)

    # Slide 2 — overall table
    fill_overall_table(slides[1], variations)

    # Slide 3 — big overall chart
    add_bar_chart(slides[2], cats, vals, "Overall conversion rate", 1.0, 1.5, 11.3, 5.4)

    # Slide 4 — device table
    fill_device_table(slides[3], variations)

    # Slide 5 — big device chart
    dcats, dvals = device_categories(variations)
    if dcats:
        add_bar_chart(slides[4], dcats, dvals, "Conversion rate by device", 1.0, 1.5, 11.3, 5.4)

    prs.save(out_path)
    print(f"Wrote {out_path} ({len(slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print(__doc__); sys.exit(1)
    with open(sys.argv[2], encoding="utf-8") as fh:
        build(sys.argv[1], json.load(fh), sys.argv[3])
