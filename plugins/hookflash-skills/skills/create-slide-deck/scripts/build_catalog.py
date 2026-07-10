"""build_catalog.py — regenerate the visual Layout Catalog for the bundled
template. Produces one slide per slide-master layout, with a red corner label
(LAYOUT <index>) and each fillable placeholder stamped with its idx + type, so
the rendered PNGs show exactly which box is which.

Run this only when the bundled template changes:
    python build_catalog.py            # -> ../references/_catalog.pptx
Then render it with render_pptx.ps1 and rebuild references/layout-catalog.md.
"""
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.normpath(os.path.join(HERE, "..", "assets", "hookflash-template.pptx"))
OUT = os.path.normpath(os.path.join(HERE, "..", "references", "_catalog.pptx"))
FURNITURE = {2, 4}  # date / slide number

prs = Presentation(TEMPLATE)

# drop any seed slides (and their relationships, so no orphan parts remain)
sldIdLst = prs.slides._sldIdLst
for s in list(sldIdLst):
    prs.part.drop_rel(s.rId)
    sldIdLst.remove(s)

master = prs.slide_masters[0]
for li, layout in enumerate(master.slide_layouts):
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx in FURNITURE:
            continue
        typ = str(ph.placeholder_format.type).split()[0]
        try:
            if ph.has_text_frame:
                ph.text_frame.text = f"[{idx}] {typ}"
        except Exception:
            pass
    tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(3400000), Emu(380000))
    tf = tb.text_frame
    tf.text = f"LAYOUT {li}"
    r = tf.paragraphs[0].runs[0]
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0x00, 0x66)

prs.save(OUT)
print("catalog deck ->", OUT, "(", len(prs.slides._sldIdLst), "layouts )")
