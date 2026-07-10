"""decklib — the deck-assembly engine for the create-slide-deck skill.

It hides python-pptx boilerplate so a build script can read like a storyboard:
open the Hookflash template, add slides *by layout index* (see the Layout
Catalog), and drop content into placeholders *by idx*. The template's slide
master carries the brand — fonts, colours, the triangle motif, dark/light
backgrounds — so you never set any of that yourself. You only supply content.

Typical build script
---------------------
    from decklib import Deck

    d = Deck()                                   # opens bundled template, empty
    s = d.add(0)                                 # layout 0 = cover
    d.text(s, 0, "Q3 Experiment Programme")      # placeholder idx 0 = title
    d.text(s, 1, "Results & recommendations")    # idx 1 = subtitle

    s = d.add(25)                                # layout 25 = case-study
    d.text(s, 0, "Sticky basket CTA")
    d.bullets(s, 31, ["Low add-to-cart on mobile", "High exit on PDP"])
    d.stat(s, 37, "+18.4%", "Add-to-cart")       # a big stat callout
    d.picture(s, 14, "chart.png")

    d.save("out.pptx")

Every method prints what it did, so a build run is self-documenting and easy to
debug. Run `python decklib.py <template.pptx>` to dump each layout's fillable
placeholders (idx + type + a size hint) — handy when planning a slide.
"""
from __future__ import annotations
import os
import sys
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_TEMPLATE = os.path.normpath(os.path.join(HERE, "..", "assets", "hookflash-template.pptx"))

# Placeholder idxs that are furniture (date / slide number) — never content.
FURNITURE_IDX = {2, 4}


class Deck:
    def __init__(self, template: str | None = None):
        self.path = template or DEFAULT_TEMPLATE
        self.prs = Presentation(self.path)
        self._layouts = list(self.prs.slide_masters[0].slide_layouts)
        # Start from a clean file: drop any seed/example slides the template ships with.
        self._clear_slides()
        print(f"[deck] opened template: {os.path.basename(self.path)} "
              f"({len(self._layouts)} layouts, cleared to 0 slides)")

    # -- internals -----------------------------------------------------------
    def _clear_slides(self):
        sldIdLst = self.prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            # drop_rel removes the relationship so the slide part is not re-serialised
            self.prs.part.drop_rel(sldId.rId)
            sldIdLst.remove(sldId)

    def _ph(self, slide, idx):
        try:
            return slide.placeholders[idx]
        except KeyError:
            avail = ", ".join(str(p.placeholder_format.idx) for p in slide.placeholders)
            raise KeyError(
                f"layout has no placeholder idx={idx}. Available idxs: [{avail}]")

    # -- building blocks -----------------------------------------------------
    def add(self, layout_index: int):
        """Add a slide using the given slide-master layout index (see catalog)."""
        layout = self._layouts[layout_index]
        slide = self.prs.slides.add_slide(layout)
        n = len(self.prs.slides.__iter__.__self__._sldIdLst)
        print(f"[deck] + slide {n} from layout {layout_index} ({layout.name!r})")
        return slide

    def text(self, slide, idx, value, *, size=None, bold=None, align=None, color=None):
        """Set a placeholder's text. Formatting inherits from the template's
        master unless you override size/bold/align/color (use sparingly — the
        template is already on-brand)."""
        ph = self._ph(slide, idx)
        tf = ph.text_frame
        tf.text = "" if value is None else str(value)
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                           "right": PP_ALIGN.RIGHT}[align]
        run = p.runs[0] if p.runs else p.add_run()
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
        return ph

    def bullets(self, slide, idx, items, *, size=None):
        """Fill a body placeholder with one paragraph per item (bullet styling
        comes from the layout)."""
        ph = self._ph(slide, idx)
        tf = ph.text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
            if size is not None:
                for r in p.runs:
                    r.font.size = Pt(size)
        return ph

    def stat(self, slide, idx, big, small=None):
        """Fill a 'stat callout' placeholder: a large number/percentage, with an
        optional smaller label on a second line. Used on the results layouts."""
        ph = self._ph(slide, idx)
        tf = ph.text_frame
        tf.clear()
        tf.paragraphs[0].text = str(big)
        if small is not None:
            p = tf.add_paragraph()
            p.text = str(small)
        return ph

    def picture(self, slide, idx, image_path):
        """Insert an image into a PICTURE placeholder (it crops-to-fill the box)."""
        ph = self._ph(slide, idx)
        ph.insert_picture(image_path)
        print(f"[deck]   picture -> idx {idx}: {os.path.basename(image_path)}")
        return ph

    def geometry(self, slide, idx):
        """Return (left, top, width, height) EMU of a placeholder — use to place
        a native chart/table exactly where a picture placeholder sits."""
        ph = self._ph(slide, idx)
        return ph.left, ph.top, ph.width, ph.height

    def remove_ph(self, slide, idx):
        """Remove a placeholder (e.g. after reading its geometry to drop a chart
        in its place)."""
        ph = self._ph(slide, idx)
        ph._element.getparent().remove(ph._element)

    def save(self, out_path):
        self.prs.save(out_path)
        n = len(self.prs.slides.__iter__.__self__._sldIdLst)
        print(f"[deck] saved {n} slides -> {out_path}")
        return out_path


def _dump(template):
    """Print each layout's fillable placeholders — planning aid."""
    prs = Presentation(template)
    for li, layout in enumerate(prs.slide_masters[0].slide_layouts):
        rows = []
        for ph in layout.placeholders:
            idx = ph.placeholder_format.idx
            if idx in FURNITURE_IDX:
                continue
            typ = str(ph.placeholder_format.type).split()[0]
            w = Emu(ph.width).inches if ph.width else 0
            h = Emu(ph.height).inches if ph.height else 0
            rows.append(f"idx {idx:>2} {typ:<12} {w:.1f}x{h:.1f}in")
        print(f"\n[{li:2}] {layout.name!r}")
        for r in rows:
            print("     ", r)


if __name__ == "__main__":
    _dump(sys.argv[1] if len(sys.argv) > 1 else DEFAULT_TEMPLATE)
